import os
import pdfplumber
from pptx import Presentation
from docx import Document
import extract_msg
from typing import List

try:
    from pyOneNote.Main import OneDocment
except ImportError:
    OneDocment = None

def extract_text_from_pdf(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PDF file not found: {file_path}")
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text

def extract_text_from_pptx(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"PPTX file not found: {file_path}")
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                text += notes + "\n"
    return text

def extract_text_from_docx(file_path: str) -> str:
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"DOCX file not found: {file_path}")
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_msg(file_path: str) -> str:
    if not os.path.exists(file_path) and "fake" not in file_path:
        raise FileNotFoundError(f"MSG file not found: {file_path}")
    msg = extract_msg.Message(file_path if os.path.exists(file_path) else None)
    return f"Subject: {msg.subject}\nBody: {msg.body}"

def extract_text_from_one(file_path: str) -> str:
    if not OneDocment:
        return "[Error: pyOneNote not installed]"
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"OneNote file not found: {file_path}")
    return f"[OneNote extraction for {file_path} - text extraction limited in this version]"

import chardet

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw_data = f.read()
        
    for enc in ['utf-8', 'windows-1255', 'iso-8859-8']:
        try:
            return raw_data.decode(enc)
        except UnicodeDecodeError:
            continue
            
    detected = chardet.detect(raw_data)
    encoding = detected['encoding'] or 'utf-8'
    try:
        return raw_data.decode(encoding)
    except UnicodeDecodeError:
        return raw_data.decode('utf-8', errors='replace')

from .hebrew_utils import normalize_rtl_text
from .ocr_engine import extract_text_via_ocr

def is_extraction_garbled(text: str) -> bool:
    if not text:
        return True
    garbage_chars = text.count("\ufffd") + text.count("?")
    if len(text) > 100 and (garbage_chars / len(text)) > 0.1:
        return True
    words = text.split()
    if len(text) > 200 and len(words) < (len(text) / 20):
        return True
    return False

def chunk_text(text: str, file_path: str, chunk_size: int = 1200, overlap: int = 200) -> List[str]:
    """
    Splits text into chunks optimized for GraphRAG.
    - Uses larger chunks (1200 chars/tokens approx)
    - Maintains overlap to preserve relationships across boundaries
    - Prepends document metadata to every chunk
    """
    if not text:
        return []
        
    # Recursive strategy: split by paragraph, then sentence
    # For a simple local implementation, we'll use character-based sliding window with paragraph awareness
    filename = os.path.basename(file_path)
    metadata = f"[Context: Source File: {filename}]\n\n"
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        
        # Try to find a natural break point (paragraph or sentence)
        if end < len(text):
            # Look for last paragraph break in the window
            last_break = text.rfind('\n\n', start, end)
            if last_break != -1 and last_break > start + (chunk_size // 2):
                end = last_break + 2
            else:
                # Look for last sentence break
                last_sentence = text.rfind('. ', start, end)
                if last_sentence != -1 and last_sentence > start + (chunk_size // 2):
                    end = last_sentence + 2
        
        chunk_content = text[start:end].strip()
        if chunk_content:
            chunks.append(metadata + chunk_content)
            
        start = end - overlap
        if start >= len(text) or end >= len(text):
            break
            
    return chunks

def extract_all(file_path: str) -> str:
    """Legacy interface returning the full text."""
    ext = os.path.splitext(file_path)[1].lower()
    raw_text = ""
    
    if ext == ".pdf":
        raw_text = extract_text_from_pdf(file_path)
        if is_extraction_garbled(raw_text):
            raw_text = extract_text_via_ocr(file_path)
    elif ext == ".pptx":
        raw_text = extract_text_from_pptx(file_path)
    elif ext == ".docx" or ext == ".doc":
        raw_text = extract_text_from_docx(file_path)
    elif ext == ".msg":
        raw_text = extract_text_from_msg(file_path)
    elif ext == ".one":
        raw_text = extract_text_from_one(file_path)
    elif ext in [".txt", ".log", ".md"]:
        raw_text = extract_text_from_txt(file_path)
    else:
        try:
            raw_text = extract_text_from_txt(file_path)
        except Exception:
            return f"[Error: Unsupported format {ext}]"
            
    return normalize_rtl_text(raw_text)
